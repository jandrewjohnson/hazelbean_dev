import os, sys, json, math, random
from collections import OrderedDict
import logging
import numpy as np
import pandas as pd
import csv
import xlrd
from PIL import Image
import io
from pptx import Presentation
import hazelbean as hb

# import numdal as nd
# import hazelbean as hb
# from hazelbean.data_structures


goals = """
Build robust pobject(dict, list, string, odict)-csv-json-markdown-yaml-xml converter,
Use recursive function to make generate_example_nested_odict complex at deeper levels.
"""

#index_synonyms = config.index_synonyms
index_synonyms = ['', 'name', 'names', 'unique_name', 'unique_names', 'index', 'indices', 'id', 'ids', 'var_name', 'var_names']

def xls_to_csv(xls_uri, csv_uri, xls_worksheet=None):
    wb = xlrd.open_workbook(xls_uri)
    if xls_worksheet:
        if isinstance(xls_worksheet, str):
            sh = wb.sheet_by_name(xls_worksheet)
        elif isinstance(xls_worksheet, int) or isinstance(xls_worksheet, float):
            sh = wb.sheet_by_index(xls_worksheet)
        else:
            raise NameError("file_to_python_object given unimplemented xls worksheet type")
    else:
        # Assume it's just the first sheet
        sh = wb.sheet_by_index(0)
    csv_file = open(csv_uri, 'w', newline='') # Python 2 version had 'wb' to avoid extra lines written. see http://stackoverflow.com/questions/3348460/csv-file-written-with-python-has-blank-lines-between-each-row
    wr = csv.writer(csv_file, quoting=csv.QUOTE_NONE, escapechar='\\') #  quoting=csv.QUOTE_ALL
    for rownum in range(sh.nrows):
        wr.writerow(sh.row_values(rownum))
    csv_file.close()


def xlsx_to_numpy_array(input_path, worksheet_name=None, skip_cols=1, skip_rows=1):
    wb = xlrd.open_workbook(input_path)
    if worksheet_name:
        if isinstance(worksheet_name, str):
            sh = wb.sheet_by_name(worksheet_name)
        elif isinstance(worksheet_name, int) or isinstance(worksheet_name, float):
            sh = wb.sheet_by_index(worksheet_name)
        else:
            raise NameError("file_to_python_object given unimplemented xls worksheet type")
    else:
        # Assume it's just the first sheet
        sh = wb.sheet_by_index(0)

    output_array = np.zeros((sh.nrows - skip_rows, sh.ncols - skip_cols))
    for rownum in range(sh.nrows - skip_rows):
        output_array[rownum] = sh.row_values(rownum + skip_rows)[skip_cols:]

    return output_array


def crop_csv_to_rect(csv_uri, data_rect):
    # Data rect is [ul row, ul col, n_rows, n_cols]
    # -1 means no limit
    for n,i in enumerate(data_rect):
        if i == -1:
            data_rect[n] = hb.config.MAX_IN_MEMORY_ARRAY_SIZE

    new_rows = []
    with open(csv_uri, 'r', newline='') as f:

        reader = csv.reader(f)
        row_id = 0
        for row in reader:
            if data_rect[0] <= row_id <= data_rect[0] + data_rect[2]:
                new_row = row[data_rect[1]: data_rect[1] + data_rect[3]]
                new_rows.append(new_row)
            row_id += 1

    with open(csv_uri, 'w', newline='') as f:
        # Overwrite the old file with the modified rows
        writer = csv.writer(f)
        writer.writerows(new_rows)


def get_strings_between_values(input_string, value_1, value_2, return_only_first=True):
    """
    Get 
    :param input_string: 
    :param value_1: 
    :param value_2: 
    :return: 
    """

    strings = []

    r = input_string.split(value_1, 1)
    if return_only_first:
        if len(r) > 1:
            r2 = r[1].split(value_2, 1)
            if len(r2) > 1:
                strings.append(r2[0])
                input_string = r2[1]
        if len(strings) > 0:
            return strings[0]
        else:
            return ''
    else:
        print ('NYI')


    



def file_to_python_object(file_uri, declare_type=None, verbose=False, return_all_parts=False, xls_worksheet=None, output_key_data_type=None, output_value_data_type=None, add_first_col_as_named_var=True):
    """
    Version that follows the simple rule of if the UL cell is blank its a DD. Else LL

    """
    if not output_value_data_type:
        output_value_data_type = str

    if not output_key_data_type:
        output_key_data_type = str

    def cast_type(input, data_type):
        if data_type is str:
            return str(input)
        elif data_type is int:
            try:
                return int(input)
            except:
                return int(float(input))
        elif data_type is float:
            return float(input)

    file_extension = None

    if os.path.exists(file_uri):
        (file_path, file_extension) = os.path.splitext(file_uri)
        (folder, file_name) = os.path.split(file_path)
    else:
        raise NameError('File given to file_to_python_object does not exist: ' + file_uri)

    if file_extension == '.json':
        json_data=open(file_uri).read()
        data = json.loads(json_data)
        return data

    elif file_extension == '.xls' or file_extension == '.xlsx':
        # If XLS, convert to a temporary csv.
        tmp_csv_uri = os.path.join(folder, file_name + '_tmp_' + hb.pretty_time() + '.csv')
        hb.remove_uri_at_exit(tmp_csv_uri)
        xls_to_csv(file_uri, tmp_csv_uri, xls_worksheet=xls_worksheet)
        file_uri = tmp_csv_uri

    data_type, num_rows, num_cols = determine_data_type_and_dimensions_from_uri(file_uri)

    if declare_type:
        data_type = declare_type

    row_headers = []
    col_headers = []

    data = None
    if data_type == 'singleton':
        with open(file_uri, 'r') as f:
            for row in f:
                split_row = row.replace('\n', '').split(',')
        data = cast_type(split_row[0], output_value_data_type)
    elif data_type == 'L':
        data = []
        with open(file_uri, 'r') as f:
            for row in f:
                split_row = row.replace('\n','').split(',')
                data.append(cast_type(split_row[0], output_value_data_type))
    elif data_type == '1d_list_vertical':
        data = []
        with open(file_uri, 'r') as f:
            for row in f:
                row = row.replace('\n','')
                data.append(cast_type(row, output_value_data_type))
    elif data_type == 'D':
        data = OrderedDict()
        with open(file_uri, 'r') as f:
            for row in f:
                split_row = row.replace('\n','').split(',')
                data[cast_type(split_row[0], output_key_data_type)] = cast_type(split_row[1], output_value_data_type)
    elif data_type == 'DD':
        data = OrderedDict()
        first_row = True
        with open(file_uri, 'r') as f:
            for row in f:
                split_row = row.replace('\n','').split(',')
                if first_row:
                    if add_first_col_as_named_var:
                        col_headers = [cast_type(i, output_key_data_type) for i in split_row]
                    else:
                        col_headers = [cast_type(i, output_key_data_type) for i in split_row[1:]]
                    first_row = False
                else:
                    row_odict = OrderedDict()
                    row_headers.append(cast_type(split_row[0], output_key_data_type))
                    for col_header_index in range(len(col_headers)):
                        if add_first_col_as_named_var:
                            row_odict[col_headers[col_header_index]] = cast_type(split_row[col_header_index], output_value_data_type)
                        else:
                            row_odict[col_headers[col_header_index]] = cast_type(split_row[col_header_index + 1], output_value_data_type) # Plus 1 because the first in the split_row is the row_header
                    data[cast_type(split_row[0], output_key_data_type)] = row_odict
    elif data_type == 'LL' or data_type == '2d_list':
        data = []
        blank_ul = False
        first_row = True
        with open(file_uri, 'r') as f:
            for row in f:
                if first_row:
                    if row.split(',')[0] in index_synonyms:
                        blank_ul = True

                if blank_ul:
                    if not first_row:
                        split_row = row.replace('\n', '').split(',')
                        data.append([cast_type(i, output_value_data_type) for i in split_row[1:]])
                else:
                    split_row = row.replace('\n', '').split(',')
                    data.append([cast_type(i, output_value_data_type) for i in split_row])

                first_row = False

    else:
        raise NameError('Unable to load file ' + file_uri + ' because datatype could not be determined from the file contents.')
    metadata = OrderedDict()
    metadata.update({'data_type':data_type,'num_rows':num_rows, 'num_cols':num_cols, 'row_headers':row_headers, 'col_headers':col_headers})

    if verbose:
        print ('\nReading file at ' + file_uri)
        print ('data_type: ' + data_type + ',  shape: num_rows ' + str(num_rows) + ', num_cols ' + str(num_cols))
        print ('col_headers: ' + ', '.join(col_headers))
        print ('row_headers: ' + ', '.join(row_headers))
        print ('python object loaded (next line):')
        print (data)

    if return_all_parts:
        return data, metadata
    else:
        return data
def save_string_as_file(input_string, file_path):
    with open(file_path, 'w') as file:
        file.write(input_string)

def determine_data_type_and_dimensions_from_uri(file_uri):
    """
    Inspects a file of type to determine what the dimensions of the data are and make a guess at the best file_type to
    express the data as. The prediction is based on what content is in the upper-left cell and the dimensions.
    Useful when converting a python iterable to a file output.
    Function forked from original found in geoecon_utils library, used with permission open BSD from Justin Johnson.

    If you dont know the file formatting types, it may be easiest to just input L, D, LL or DD for List, Dictionary, 2dim List,
    2dmin Dicitonary manually

    :param file_uri:
    :return: data_type, num_rows, num_cols
    """

    row_headers = []
    col_headers = []

    if os.path.exists(file_uri):
        # Iterate trough one initial time to ddetermine dimensions and
        blank_ul = False
        with open(file_uri, 'r') as f:
            # Save all col_lengths to allow for truncated rows.
            col_lengths = []
            first_row = True
            for row in f:
                split_row = row.replace('\n', '').split(',')
                col_lengths.append(len(split_row))
                if first_row:
                    if split_row[0] in index_synonyms:
                        blank_ul = True
                    else:
                        blank_ul = False


        num_rows = len(col_lengths)
        num_cols = max(col_lengths)

        # with open(file_uri, 'r') as f:
        data_type = 'empty'
        if num_cols == 1:
            if num_rows == 1:
                if blank_ul:
                    data_type = 'empty'
                else:
                    data_type = 'singleton'
            else:
                if blank_ul:
                    data_type = 'row_headers'
                else:
                    data_type = 'horizontal_list'
        elif num_cols >= 2:
            if num_rows == 1:
                if blank_ul:
                    data_type = 'col_headers'
                else:
                    data_type = 'horizontal_list'
            if num_rows >= 2:
                if blank_ul:
                    data_type = 'DD'
                else:
                    data_type = 'LL'

        return data_type, num_rows, num_cols
    else:
        raise NameError('File given to ' + file_uri + ' determine_data_type_and_dimensions_for_read does not exist.')


def determine_data_type_and_dimensions_from_object(input_python_object):
    """
    Inspects a file of type to determine what the dimensions of the data are and make a guess at the best file_type to
    express the data as. The prediction is based on what content is in the upper-left cell and the dimensions.
    Useful when converting a python iterable to a file output.
    Function forked from original found in geoecon_utils library, used with permission open BSD from Justin Johnson.
    """

    data_type = None
    blank_ul = False

    # First check to see if more than 2 dimensions. Currently, I do not detect beyond 2 dimensions here and instead just use the
    # Str function in python in the write function.
    if isinstance(input_python_object, str):
        data_type = 'singleton'
    elif isinstance(input_python_object, dict):
        raise TypeError('Only works with OrderedDicts not dicts.')
    elif isinstance(input_python_object, list):
        first_row = input_python_object[0]
        if isinstance(first_row, (str, int, float, bool)):
            data_type = 'vertical_list'
        elif isinstance(first_row, dict):
            raise TypeError('Only works with OrderedDicts not dicts.')
        elif isinstance(first_row, list):
            if first_row[0]: # If it's blank, assume it's an empty headers column
                data_type = 'LL'
            else:
                data_type = 'column_headers'
        elif isinstance(first_row, OrderedDict):
            data_type = 'LD' # Unimplemented
        else:
            raise NameError('type unknown')
    elif isinstance(input_python_object, OrderedDict):
        first_row_key = next(iter(input_python_object))
        first_row = input_python_object[first_row_key]
        if isinstance(first_row, (str, int, float, bool)):
            data_type = 'D'
        elif isinstance(first_row, dict):
            raise TypeError('Only works with OrderedDicts not dicts.')
        elif isinstance(first_row, list):
            data_type = 'DL' # NYI
            raise TypeError(data_type + ' unsupported.')
        elif isinstance(first_row, OrderedDict):
            data_type = 'DD'
        else:
            raise NameError('Unsupported object type. Did you give a blank OrderedDict to python_object_to_csv()?. \nYou Gave:\n\n' + str(input_python_object))
    else:
        raise NameError('Unsupported object type. You probably gave "None" to python_object_to_csv()')
    return data_type

def dict_to_df(input_dict):
    """Fast way to go from double-named nested dictionary to DF. Assumes row-column orientation.
    Example: input_dict[row_label][col_label] = Dict
        """
    data_dict = {}
    columns = None
    for k1 in input_dict.keys():
        data_dict[k1] = []
        for k2 in input_dict[k1].keys():
            if columns is None:
                columns = list(input_dict[k1].keys())
            data_dict[k1].append(input_dict[k1][k2])

    return pd.DataFrame.from_dict(data_dict, orient='index', columns=columns)

def df_to_dict(input_df):
    """Fast way to go from DF to double-named nested dictionary. Assumes row-column orientation.
    Example: input_dict[row_label][col_label] = Dict
        """
    output_dict = {}
    indices = list(input_df.index)
    columns = list(input_df.columns)
    for index in indices:
        output_dict[index] = {}
        for column in columns:
            output_dict[index][column] = input_df.loc[index].loc[column]

    return output_dict


def python_object_to_csv(input_iterable, output_uri, csv_type=None, verbose=False):
    """Improved by dict_to_df()"""
    if csv_type:
        data_type = csv_type
    else:
        data_type = determine_data_type_and_dimensions_for_write(input_iterable)
    protected_characters = [',', '\n']
    first_row = True

    if not os.path.exists(os.path.split(output_uri)[0]) and os.path.split(output_uri)[0]:
        print (('Specified output_uri folder ' + os.path.split(output_uri)[0] + ' does not exist. Creating it.'))
        os.makedirs(os.path.split(output_uri)[0])

    to_write = ''
    if data_type == 'singleton':
        to_write = input_iterable

    # TODO This is a potential simpler way: just declare it.
    elif data_type == 'rc_2d_odict':
        first_row = True
        for key, value in list(input_iterable.items()):
            if first_row:
                first_row = False
                to_write += ',' + ','.join(list(value.keys())) + '\n'

            write_list = [str(i) for i in list(value.values())]
            to_write += str(key)+ ',' + ','.join(write_list) + '\n'

    elif data_type == 'cr_2d_odict':
        first_row = True
        for key, value in list(input_iterable.items()):
            if first_row:
                first_row = False
                to_write += ',' + ','.join(list(value.keys())) + '\n'

            write_list = [str(i) for i in list(value.values())]
            to_write += str(key)+ ',' + ','.join(write_list) + '\n'


    elif data_type == '1d_list_vertical':
        to_write += '\n'.join(input_iterable)
    elif data_type == '1d_list_horizontal':
        to_write += ','.join(input_iterable)
    elif data_type == '1d_dict_vertical':
        to_write += '\n'.join(str(k) + ',' + str(v) for k, v in input_iterable.items())

    elif data_type == '2d_list':
        for row in input_iterable:
            if any(character in row for character in protected_characters):
                raise NameError('Protected character found in the string-ed version of the iterable.')

            if len(row) > 0:
                if not isinstance(row[0], str):
                    row = [str(i) for i in row]
                to_write += ','.join(row) + '\n'
    elif data_type == '2d_cr_list':
        for row in input_iterable:
            if any(character in row for character in protected_characters):
                raise NameError('Protected character found in the string-ed version of the iterable.')

            if len(row) > 0:
                if not isinstance(row[0], str):
                    row = [str(i) for i in row]
                to_write += ','.join(row) + '\n'
    elif data_type == '2d_rc_list':

        for row in input_iterable:
            if any(character in row for character in protected_characters):
                raise NameError('Protected character found in the string-ed version of the iterable.')

            if len(row) > 0:
                if not isinstance(row[0], str):
                    row = [str(i) for i in row]
                to_write += ','.join(row) + '\n'
    elif data_type == '2d_list_odict_NOT_SUPPORTED':
        raise NameError('2d_list_odict_NOT_SUPPORTED unknown')
    elif data_type == '1d_odict':
        for key, value in list(input_iterable.items()):
            # check to see if commas or line breaks are in the iterable string.
            value = str(value)
            if any(character in str(key) for character in protected_characters) or any(character in value for character in protected_characters):
                raise NameError('Protected character found in the string-ed version of the iterable: '+ str(key))
            to_write += str(key) + ',' + str(value) + '\n'
    elif data_type == '2d_odict_list':
        lenghts = [len(v) for k, v in input_iterable.items()]
        max_length = max(lenghts)
        new_dict = {}
        for k, v in input_iterable.items():
            if len(v) < max_length:
                new_dict[k] = v + [''] * (max_length - len(v))
            else:
                new_dict[k] = v
        df = pd.DataFrame(data=new_dict)
        df.to_csv(output_uri, index=False)
        
        return
                
            
        # raise NameError('2d_odict_list unknown')
    elif data_type == '2d_odict':
        if isinstance(input_iterable, list):
            # The only way you can get here is it was manually declared to be this type and the list implies that it was empty (1 row).
            # TODOO Currently, I do not deal with indexed data_types consistently, nor do I account for empty data (as in here) the same on IO operations.
            to_write += ','.join(input_iterable)
        else:
            for key, value in list(input_iterable.items()):
                if first_row:
                    # On the first row, we need to write BOTH th efirst and second rows for col_headers and data respecitvely.
                    if key is not None and value is not None:
                        if any(character in str(key) for character in protected_characters) or any(character in value for character in protected_characters):
                            raise NameError('Protected character found in the string-ed version of the iterable.')
                    to_write += ','.join([''] + [str(i) for i in value.keys()]) + '\n' # Note the following duplication of keys, values to address the nature of first row being keys.
                if key is not None and value is not None:
                    if any(character in str(key) for character in protected_characters) or any(character in value for character in protected_characters):
                        raise NameError('Protected character found in the string-ed version of the iterable.')
                    first_col = True
                    for value2 in list(value.values()):
                        if first_col:
                            # if first_row:
                            #     to_write += ','
                            to_write += str(key) + ','
                            first_col = False
                        else:
                            to_write += ','
                        if isinstance(value2, list):
                            to_write += '<^>'.join(value2)
                        else:
                            to_write += str(value2)
                    to_write += '\n'
                    first_row = False
                else:
                    to_write +=','
    else:
        raise NameError('Not sure how to handle that data_type.')


    # Write the string to a file
    with open(output_uri, 'w', newline='', encoding='utf-8') as file:
        file.write(to_write)        
    # with open('n_cells_per_zone.csv', 'w', newline='') as file:
    #     writer = csv.writer(file)
    #     writer.writerows(data)
    #     open(output_uri, 'w', encoding='utf').write(to_write)

    if verbose:
        print (('\nWriting python object to csv at ' + output_uri + '. Auto-detected the data_type to be: ' + data_type))
        print (('String written:\n' + to_write))


def determine_data_type_and_dimensions_for_write(input_python_object):
    """
    Inspects a file of type to determine what the dimensions of the data are and make a guess at the best file_type to
    express the data as. The prediction is based on what content is in the upper-left cell and the dimensions.
    Useful when converting a python iterable to a file output.
    Function forked from original found in geoecon_utils library, used with permission open BSD from Justin Johnson.
    """

    data_type = None

    # First check to see if more than 2 dimensions. Currently, I do not detect beyond 2 dimensions here and instead just use the
    # Str function in python in the write function.
    if isinstance(input_python_object, str):
        data_type = 'singleton'
    # elif isinstance(input_python_object, dict):
    #     raise TypeError('Only works with OrderedDicts not dicts.')
    elif isinstance(input_python_object, list):
        first_row = input_python_object[0]
        if isinstance(first_row, (str, int, float, bool)):
            data_type = '1d_list'
        elif isinstance(first_row, dict):
            raise TypeError('Only works with OrderedDicts not dicts.')
        elif isinstance(first_row, list):
            data_type = '2d_list'
        elif isinstance(first_row, OrderedDict):
            data_type = '2d_list_odict_NOT_SUPPORTED'
        else:
            raise NameError('2d_list_odict_NOT_SUPPORTED unknown')
    elif isinstance(input_python_object, OrderedDict) or isinstance(input_python_object, dict):
        first_row_key = next(iter(input_python_object))
        first_row = input_python_object[first_row_key]
        if isinstance(first_row, (str, int, float, bool, np.ubyte, np.byte, np.int16, np.int32, np.int64, np.float32, np.float64)):
            data_type = '1d_odict'
        # elif isinstance(first_row, dict):
        #     raise TypeError('Only works with OrderedDicts not dicts.')
        elif isinstance(first_row, list):
            data_type = '2d_odict_list'
        elif isinstance(first_row, OrderedDict) or isinstance(first_row, dict):
            data_type = '2d_odict'
        else:
            raise NameError('Unsupported object type. Did you give a blank OrderedDict to python_object_to_csv()?. \nYou Gave:\n\n' + str(input_python_object))
    elif isinstance(input_python_object, np.ndarray):
        data_type = '2d_list'

    else:
        raise NameError('Unsupported object type. You probably gave "None" to python_object_to_csv()')
    return data_type

def comma_linebreak_string_to_2d_array(input_string, dtype=None):
    """Only good for small arrays for testing purposes. Not vectorized."""
    s = str(input_string)
    rows = s.split('\n')

    # First get size
    n_rows = len(rows)
    n_cols = len(rows[0].split(','))

    if dtype:
        a = np.zeros((n_rows, n_cols), dtype=dtype)
    else:
        a = np.zeros((n_rows, n_cols))

    for row_id, row in enumerate(rows):
        col = row.split(',')
        for col_id, value in enumerate(col):
            a[row_id][col_id] = value

    return a

def dictionary_to_dataframe(input_dictionary, output_path=None):
    """Input dictionary or OrderedDict should be 2-dimension in row-column nesting."""
    if type(input_dictionary) not in [OrderedDict, dict]:
        raise NameError('dictionary_to_dataframe only works with dicts or OrderedDicts. You gave: ' + str(input_dictionary))

    # Inspect first entry in dict to get column names
    extract_columns = []
    for k, v in input_dictionary.items():
        for kk, vv in v.items():
            extract_columns.append(kk)
        break

    output_df = pd.DataFrame.from_dict(input_dictionary, orient='index', columns=extract_columns)
    if output_path is not None:
        if os.path.splitext(output_path)[1] == '.csv':
            output_df.to_csv(output_path)
        elif os.path.splitext(output_path)[1] in ['.xls', '.xlsx']:
            output_df.to_excel(output_path)
    return output_df

# def read_csv(input_path):
#     df = pd.read_csv(input_path)
    
#     for col in df.columns:
#         if col.split('_')[-1] == 'list':
#             for 

def propose_fuzzy_merge(left_df, right_df, on=None, left_on=None, right_on=None, how=None,
                        fuzzy_merge_csv_path=None, fuzzy_merge_report_path=None, cutoff=0.65):

    if on is not None:
        left_on = on
        right_on = on

    left_uniques = pd.unique(left_df[left_on])
    right_uniques = pd.unique(right_df[right_on])

    # comparison_dict = hb.compare_sets_as_dict(left_uniques, right_uniques)


    potential_merge, report = hb.fuzzy_merge(left_df, right_df, left_on, right_on, how=how, cutoff=cutoff)


    # Write the report dict to a csv if path given
    report_string = "remap_input, remap_output, option_2, option_3, option_4, option_5, option_6, option_7\n"
    for k, v in report.items():
        if v is not None:
            if len(v) > 0:
                if len(v[1][1:]) > 7:
                    last_index = 7
                else:
                    last_index = len(v[1][1:])

                report_string += "\"" + str(k) + "\",\"" + str(v[0]) + "\"," + ",".join(["\"" + str(i) + "\"" for i in v[1][1:last_index]])
            else:
                report_string += "\"" + str(k) + "\""
        else:
            report_string += "\"" + str(k) + "\""
        report_string += "\n"


    from io import StringIO
    df_string = StringIO(report_string)
    report = pd.read_csv(df_string)

    potential_left_uniques = pd.unique(left_df[left_on])
    potential_right_uniques = pd.unique(right_df[right_on])
    potential_merge_uniques = pd.unique(potential_merge[left_on])
    potential_merge_uniques = {i for i in potential_merge_uniques if i is not None}

    report.sort_values('remap_input', inplace=True)
    report.reset_index(inplace=True)

    report = pd.concat([report, pd.Series(sorted(left_uniques), name='possible_labels_from_left')], axis=1)
    report = pd.concat([report, pd.Series(sorted(right_uniques), name='possible_labels_from_right')], axis=1)
    # report['possible_labels_from_left'] = pd.Series(left_uniques)
    # report['possible_labels_from_right'] = pd.Series(right_uniques)


    # report['possible_labels_from_left'] = report['possible_labels_from_left'].sort_values()
    # report['possible_labels_from_right'] = report['possible_labels_from_right'].sort_values()
    report = report[['possible_labels_from_left', 'possible_labels_from_right'] + [i for i in report.columns if i not in ['possible_labels_from_left', 'possible_labels_from_right', 'index']]]


    # At this point, left_uniques should be the same as merge uniques (cause we assigned right into left)
    comparison_dict_left = hb.compare_sets_as_dict(potential_left_uniques, potential_merge_uniques)


    comparison_dict_right = hb.compare_sets_as_dict(potential_right_uniques, potential_merge_uniques)
    

    # pd.merge(food_balance_pivot_reduced_df, food_balance_pivot_reduced_df, on='item_name', how='outer')

    if fuzzy_merge_report_path is not None:
        report.to_csv(fuzzy_merge_report_path, index=False)
    else:
        pass

    if fuzzy_merge_csv_path is not None:
        potential_merge.to_csv(fuzzy_merge_csv_path)
    return potential_merge, report, comparison_dict_left, comparison_dict_right

def merge_dataframes_with_remap(left_df, right_df, remap_df_or_path, remap_left_col='remap_input', remap_right_col='remap_output', on=None, left_on=None, right_on=None, how='outer'):
    if isinstance(remap_df_or_path, str):
        if remap_df_or_path.endswith('.csv'):
            remap_df_or_path = pd.read_csv(remap_df_or_path, encoding='latin')
        else:
            remap_df_or_path = pd.read_excel(remap_df_or_path)

    if on is not None:
        if left_on is not None:
            raise NameError('Cannot give both on and left_on')
        if right_on is not None:
            raise NameError('Cannot give both on and right_on')
        left_on = on
        right_on = on


    left_values = remap_df_or_path[remap_left_col]

    right_values = remap_df_or_path[remap_right_col]
    remap_dict = dict(zip(left_values, right_values))


    right_df[right_on] = right_df[right_on].replace(remap_dict)

    # remap_dict = {k: v for k, v in dict(zip(left_values, right_values).items()}

    merged_df = pd.merge(left_df, right_df, left_on=left_on, right_on=right_on, how=how)
    return merged_df

def check_if_has_key(notebook_path):
    """Bespoke for quarto notebooks and related publishing. Checks if the notebook has a key in the first cell."""
    with open(notebook_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        has_key = False
        for v in data['cells']:
            if v['source'][0].lstrip().startswith('# KEY'):
                has_key = True
                break
    return has_key

def strip_quarto_header_from_ipynb(input_path, output_path):
    
    """Bespoke for quarto notebooks and related publishing. Removes the quarto header from an ipynb file.
    """
    skip_block = False
    to_write = []
    # print('hb.strip_quarto_header_and_keys_from_ipynb() on ', input_path, output_path)
    with open(input_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        output_dict = {}
        output_dict['cells'] = []
        for v in data['cells']:
            if not v['source'][0].lstrip().startswith('---'):
                output_dict['cells'].append(v)
                
                
    output_dict['metadata'] = {
            "kernelspec": {
            "name": "python3",
            "language": "python",
            "display_name": "Python 3 (ipykernel)"
            }        },  
    output_dict['nbformat'] = 4
    output_dict['nbformat_minor'] = 4
    # print('output_path', output_path, '\n', output_dict)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(output_dict, f, ensure_ascii=False, indent=4) 
          
def strip_quarto_header_and_keys_from_ipynb(input_path, output_path):
    """Bespoke for quarto notebooks and related publishing. Removes the quarto header from an ipynb file.
    """

    # print('hb.strip_quarto_header_and_keys_from_ipynb() on ', input_path, output_path)
    with open(input_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        output_dict = {}
        output_dict['cells'] = []
        
        for v in data['cells']:
            if not v['source'][0].lstrip().startswith('---') and not v['source'][0].lstrip().startswith('# KEY'):
                output_dict['cells'].append(v)
                
                
    output_dict['metadata'] = {
            "kernelspec": {
            "name": "python3",
            "language": "python",
            "display_name": "Python 3 (ipykernel)"
            }        },  
    output_dict['nbformat'] = 4
    output_dict['nbformat_minor'] = 4
    # print('output_path', output_path, '\n', output_dict)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(output_dict, f, ensure_ascii=False, indent=4)                
    # # Write it
    # with open(output_path, 'w', encoding='utf-8') as f:
    #     f.writelines(to_write)
                    


def pptx_to_markdown_slide_dict(pptx_path, output_dir=None):
    
    if output_dir is None:
        output_dir = os.path.dirname(pptx_path)
        
    prs = Presentation(pptx_path)
    markdown_slide_dict = {}
    image_count = 0
    for slide_number, slide in enumerate(prs.slides):
        markdown_slide_dict[slide_number] = {}
        
        if slide_number == 3:
            pass
            
        actual_shape_number = 0
        for shape_number, shape in enumerate(slide.shapes):
            
            markdown_slide_dict[slide_number][actual_shape_number] = ''
            if hasattr(shape, "text_frame"):
                paragraphs = shape.text_frame.paragraphs

                for paragraph in paragraphs:
                    level = paragraph.level
                    bullet_prefix = ' ' * (4 * level) + ('- ' if len(paragraphs) > 0 and actual_shape_number > 0 else '')
                    to_append = bullet_prefix + paragraph.text + ('\n' if len(paragraphs) > 1 else '')
                    if len(paragraph.text) > 0:
                        markdown_slide_dict[slide_number][actual_shape_number] += to_append  
                if len(markdown_slide_dict[slide_number][actual_shape_number]) > 0: 
                    actual_shape_number += 1                
            
            elif shape.shape_type == 13:
                image_stream = io.BytesIO(shape.image.blob)
                hb.create_directories(f"{output_dir}/images")

                with Image.open(image_stream) as img:
                    try:
                        time = hb.pretty_time()
                        img_path = f"{output_dir}/images/image_{image_count}_{time}.png"
                        img.save(img_path)
                        image_count += 1    
                    except:
                        raise NameError('Error saving image: ', img_path)

                    img_hash = hb.hash_file_path(img_path)
                    parent_dir = os.path.split(os.path.dirname(pptx_path))[1]
                    filename = os.path.split(pptx_path)[1].replace('.', '_')
                    qmd_embed_image_path = os.path.join('images', 'image_' + parent_dir + '_' + filename + '_' + img_hash + '.png')
                    hashed_img_path = os.path.join(output_dir, qmd_embed_image_path)
                    
                    hb.path_rename(img_path, hashed_img_path, skip_if_exists=True)
                    
                    # If it still exists (cause the hasehd version already existed) delete it
                    if hb.path_exists(img_path):
                        hb.remove_path(img_path)
                        
                        
                    image_md = f"![]({qmd_embed_image_path})\n"
                    
                    # Add the markdown to display the image
                    markdown_slide_dict[slide_number][actual_shape_number] += image_md
                    actual_shape_number += 1


            
            elif shape.shape_type == 14: # PLACEHOLDER TEXT?
                if hasattr(shape, "text"): 
                    markdown_slide_dict[slide_number][actual_shape_number] += f"- {shape.text}\n"
                    actual_shape_number += 1
                else:
                    current_text = ""                
                    markdown_slide_dict[slide_number][actual_shape_number] += f"{current_text}\n"
                    actual_shape_number += 1
                
            elif shape.shape_type == 17: # TEXT BOX
                markdown_slide_dict[slide_number][actual_shape_number] += f"- {shape.text}\n"
                actual_shape_number += 1
            else:
                print('Unknown shape type', shape.shape_type)
                
    # Clean the dict
    output_dict = {}
    slide_counter = 0
    for slide_number, slide_dict in markdown_slide_dict.items():
        
        # Strip empty elements
        if len(slide_dict) > 0:
            output_dict[slide_number] = {}
            
        # Replace images with text if there is text in the same slide ## LEARNING POINT: This has to be done before you iterate through the dict or it raises a dict changes when iterated.           
        if slide_dict[0].lstrip().startswith('![') and len(slide_dict) > 1: # Then it's a slide with an image as its first element and it has something else in it.
            for shape_number, shape_text in slide_dict.items():
                if not shape_text.lstrip().startswith('!['):
                    key_to_move = shape_number
                    break
            value_to_move = slide_dict.pop(key_to_move)
            if value_to_move[0:2] == '- ': # If it starts with a bullet point, remove it
                value_to_move = value_to_move[2:]
            slide_dict = {key_to_move: value_to_move, **slide_dict}
                                           
        # Iterate through slides  
        shape_counter = 0
        for shape_number, shape_text in slide_dict.items():
            
            if len(shape_text) > 0:
                if shape_text.endswith('\n'):
                    shape_text = shape_text[:-1]
                output_dict[slide_counter][shape_counter] = shape_text
                

                        
                # # create a new dict with the found text as the first
                # new_dict = {}
                # for shape_number in range(len(slide_dict)):
                    
                #     new_dict[shape_number]
                #     output_dict[slide_counter][shape_counter] = found_text
            
            
            shape_counter += 1   
             
        slide_counter += 1
            
    return output_dict

def markdown_slide_dict_to_qmd(markdown_slide_dict, output_path, theme='sky'):

    output_string = ""
    for slide_number, slide in markdown_slide_dict.items():
        if slide_number == 3:
            pass
        for shape_number, shape in slide.items():

            if shape_number == 0:
                if len([v for k, v in slide.items() if len(v.replace('\n', ''))]) == 1:

                    output_string += "# " + shape + "\n\n"
                else:
                    output_string += "## " + shape + "\n\n"
            else:
                if len(shape.replace('\n', '')) > 0:
                    output_string += shape + '\n\n'

        revealjs_prepend = """---
format: 
    revealjs:
        theme: """ + theme + """
        margin: 0       
        self-contained: false
        scrollable: true
        code-fold: show
        slide-number: true  
        preview-links: auto
        css: styles.css
        incremental: true  
        auto-stretch: false
---

"""

    hb.write_to_file(revealjs_prepend + output_string, output_path)

# TODOOO: make it so <nonincremental> still has r-fit-text

def qmd_path_to_marked_qmd_path(qmd_path, marked_qmd_path):
    
    # TODOOO: Continue going through and adding additional tags like imgbg but for things like 2cols
    
    output_lines = []
    # Read the qmd file but with utf8
    qmd_header_read_status = 0
    new_slide = 0
    header_lines = []
    fit_threshold = 300
    
    with open(qmd_path, 'r', encoding='utf-8') as f:
        data = f.read()   
        lines = data.split('\n')             
        n_lines = len(lines)
        content_lines = []  
        
        # Get all lines in the list between the instances of '---'
        for line in lines:
            if line.startswith('---') and qmd_header_read_status == 0:
                qmd_header_read_status = 1
                header_lines.append(line)
            elif line.startswith('---') and qmd_header_read_status == 1:
                qmd_header_read_status = 0
                header_lines.append(line)
                header_lines.append('')
            elif qmd_header_read_status == 1:
                header_lines.append(line)
            else:
                content_lines.append(line)


        output_content_lines = []

        
        if len(header_lines) == 0:
            header_lines = """---
format: 
    revealjs:
        theme: simple
        margin: 0       
        self-contained: false
        scrollable: true
        code-fold: show
        slide-number: true  
        preview-links: auto
        css: styles.css
        incremental: true  
        auto-stretch: false
---""".split('\n')
        
        # Iterate through content_lines
        slide_content = []
        tags = []
        for c, line in enumerate(content_lines):
            

            if line.startswith('# '):
                # New section slide
                if line.rstrip().endswith('>'):
                    # append all instances of characters between < and > into a list

                    for i in range(len(line)):
                        if line[i] == '<':
                            start = i
                        if line[i] == '>':
                            end = i
                            tags.append(line[start:end+1])
                new_slide = 1
                
            elif line.startswith('## '):
                # New regular slide
                if line.rstrip().endswith('>'):
                    # append all instances of characters between < and > into a list
     
                    for i in range(len(line)):
                        if line[i] == '<':
                            start = i
                        if line[i] == '>':
                            end = i
                            tags.append(line[start:end+1])
                new_slide = 1
            else:
                new_slide = 0
            
            # Get all the content of that slide 
            if new_slide:
                # Get tags
                tags = []
                already_inserted_new_col = False
                slide_content = []
                for item_c in range(len(content_lines)):
                    if item_c + c + 1 >= len(content_lines):
                        break
                    if content_lines[item_c + c + 1].lstrip().startswith('#'):
                        break
                    else:
                        slide_content.append(content_lines[item_c + c])
                slide_content.append('')
                
            n_images = 0
            # Count how many instances of ![ there are in the slide_content
            for line in slide_content:
                if line.lstrip().startswith('!['):
                    n_images += 1

            if n_images >= 1:
                if new_slide:
                    tags.append('<list-left-images-right>')
                    
                    # TODOOO, this isn't getting added
                    slide_content[0] = slide_content[0] + ' <list-left-images-right>'
                
            
            # Iterate through tags to modify slide_content
            if len(tags) > 0:
                for tag in tags:
                    if tag == '<imgbg>':
                        
                        # Find the first image                        
                        for i in range(len(slide_content)):
                            if slide_content[i].lstrip().startswith('!['):
                                bg_image = slide_content[i]                                
                                del slide_content[i]       
                                break
                            
                        bg_image_filepath = bg_image.split('(')[1].split(')')[0]
                    
                        for i in range(len(slide_content)):
                            
                            if '<imgbg>' in slide_content[i]:
                                output_content_lines.append(slide_content[i] + ' {background-image="' + bg_image_filepath + '" background-size="contain" background-repeat="no-repeat" }\n')
                                
                            elif len(slide_content[i]) > 0:
                                # Add a white background with transparency. Because this is just a p, style it with a div rather than via the style.css as was done for the header.
                                output_content_lines.append("<div style=\"background-color: rgba(255, 255, 255, 0.4); color: black; padding: 20px;\">" + slide_content[i] + '</div>\n')
                    elif tag == "<nonincremental>":
                        slide_content[0] = slide_content[0].replace('<nonincremental>', "\n\n::: {.nonincremental}")
                        slide_content.insert(-1, ":::")
                        output_content_lines += slide_content
                    elif tag == "<list-left-images-right>":
                        slide_content[0] = slide_content[0].replace('<list-left-images-right>', "\n\n::: columns \n::: {.column width=\"70%\"} \n::: r-fit-text \n")
                        to_insert = []
                        for c, slide_content_line in enumerate(slide_content):
                            if slide_content_line.lstrip().startswith('![') and not already_inserted_new_col:
                                to_insert.append((c-1, '::: \n::: \n::: {.column width=\"30%\"}'))
                                already_inserted_new_col = True
                                break
                            # output_content_lines.append(slide_content_line)
                        for i, v in to_insert:
                            slide_content.insert(i, v)
                        
                        if new_slide: # NOTE AWKWARDNESS. I create the whole slide's content. 
                            #Then i iterate over each line, creating the slide's content once for each line...
                            # Would be better to only generate it once.
                            slide_content.append(':::\n:::\n')
                            output_content_lines += slide_content
            else:
                if new_slide:
                    if 'imgbg' not in tags and 'list-left-images-right' not in tags:
                        
                        if len(str(slide_content)) >= fit_threshold:
                            slide_content.insert(2, "::: r-fit-text")
                            slide_content[-1] = ':::\n'
                        
                        # AUTOSCALING NOT WORKING
                        # for content_c, content in enumerate(slide_content):
                        #     to_postpend = ''
                        #     if content.lstrip().startswith('![') and 'width' not in content:
                        #         to_postpend += " width=\"900\" "
                        #     # if content.lstrip().startswith('![') and 'height' not in content:
                        #     #     to_postpend += " height=\"30\" "     
                        #     if len(to_postpend) > 0:
                        #         slide_content[content_c] = content + '{' + to_postpend + '}'
                    output_content_lines += slide_content
                    
        

                        
            
            
    output_lines = header_lines + output_content_lines
    output_string = '\n'.join(output_lines)
    hb.write_to_file(output_string, marked_qmd_path)



    


def qmd_to_revealjs(src_qmd_path):
           
    os.system("quarto preview \"" + src_qmd_path + "\"")
    






